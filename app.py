# app.py
import os
import json
from pathlib import Path
from typing import List, Dict, Any

import streamlit as st
import pdfplumber
from pptx import Presentation
from openai import OpenAI


# Session state (persist result)

if "final" not in st.session_state:
    st.session_state.final = None
if "txt_content" not in st.session_state:
    st.session_state.txt_content = ""
if "last_outline" not in st.session_state:
    st.session_state.last_outline = None
if "scroll_to_agent" not in st.session_state:
    st.session_state.scroll_to_agent = False


# Config

MODEL = "gpt-4o-mini"
OUTPUT_SCHEMA_HINT = {
    "emails": [{"subject": "string", "body": "string", "cta": "string"}],
    "linkedin": ["string"],
    "script": "string",
}


# Utilities

def load_css(path: str) -> None:
    try:
        css = Path(path).read_text(encoding="utf-8")
        st.markdown(f"<style>{css}</style>", unsafe_allow_html=True)
    except FileNotFoundError:
        pass

def get_qs(key: str, default: str) -> str:
    try:
        return st.query_params.get(key, default)
    except Exception:
        qp = st.experimental_get_query_params()
        return qp.get(key, [default])[0]

def set_qs(**kwargs) -> None:
    try:
        st.query_params.update(kwargs)
    except Exception:
        st.experimental_set_query_params(**kwargs)

def inject_scroll_to(anchor_id: str) -> None:
    st.markdown(
        f"""
<script>
  const el = window.parent.document.getElementById("{anchor_id}");
  if (el) {{
    el.scrollIntoView({{behavior: "smooth", block: "start"}});
  }}
</script>
""",
        unsafe_allow_html=True,
    )


# File extraction

def extract_pdf_text(path: str) -> str:
    parts: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            txt = page.extract_text() or ""
            if txt.strip():
                parts.append(txt.strip())
    return "\n\n---\n\n".join(parts)

def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    slides: List[str] = []
    for slide in prs.slides:
        items: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    items.append(t)
        if items:
            slides.append("\n".join(items))
    return "\n\n---\n\n".join(slides)

def to_outline(raw: str) -> Dict[str, Any]:
    chunks = [s.strip() for s in raw.split("\n\n---\n\n") if s.strip()]
    if not chunks:
        chunks = [raw.strip()]
    slides = [{"index": i + 1, "text": s} for i, s in enumerate(chunks[:60])]
    return {"slides": slides}


# OpenAI helpers

def openai_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        st.error("Missing OPENAI_API_KEY environment variable")
        st.stop()
    return OpenAI(api_key=api_key)

def chat_json(client: OpenAI, system: str, user: str) -> str:
    resp = client.chat.completions.create(
        model=MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        temperature=0.2,
    )
    return resp.choices[0].message.content

def coerce_json(text: str) -> dict:
    t = text.strip()
    if t.startswith("```"):
        t = t.strip("`")
        nl = t.find("\n")
        if nl != -1:
            t = t[nl + 1:].strip()
    try:
        return json.loads(t)
    except Exception:
        pass
    start = t.find("{")
    end = t.rfind("}")
    if start != -1 and end != -1 and end > start:
        snippet = t[start:end + 1]
        try:
            return json.loads(snippet)
        except Exception:
            st.error("Model returned invalid JSON. Try again or simplify style notes.")
            st.stop()
    st.error("No JSON found in model response.")
    st.stop()

def plan_brief(client: OpenAI, outline: dict, persona: str, tone: str, style: str) -> dict:
    system = "You extract concise marketing briefs for outbound work. Return pure JSON."
    user = (
        f"Persona: {persona}\n"
        f"Tone: {tone}\n"
        f"Style notes: {style}\n"
        f"Deck outline JSON:\n{json.dumps(outline, ensure_ascii=False)}\n\n"
        "Return JSON with fields: audience (string), pains (array of strings), "
        "value (array of strings), proofs (array of strings, only from the deck), "
        "objections (array of strings)."
    )
    return coerce_json(chat_json(client, system, user))

def generate_assets(client: OpenAI, brief: dict) -> dict:
    system = (
        "You are a senior B2B marketer. "
        "Output pure JSON matching this schema:\n"
        f"{json.dumps(OUTPUT_SCHEMA_HINT)}"
    )
    rules = (
        "Write a 5 step outbound email sequence, 3 LinkedIn posts, and a 90 second talk track.\n"
        "Rules:\n"
        "- One CTA per email.\n"
        "- Subject under 50 characters.\n"
        "- Email body under 140 words.\n"
        "- Plain language.\n"
        "- Every claim must map to a proof in brief.proofs.\n"
        "Return JSON only with keys: emails, linkedin, script."
    )
    user = f"Brief JSON:\n{json.dumps(brief, ensure_ascii=False)}\n\n{rules}"
    return coerce_json(chat_json(client, system, user))

def guardrail_review(client: OpenAI, draft: dict, style: str, brief: dict) -> dict:
    system = (
        "You review for clarity, tone, and proof alignment. "
        "Keep the same JSON keys and structure. Return pure JSON."
    )
    user = (
        f"Draft JSON:\n{json.dumps(draft, ensure_ascii=False)}\n\n"
        f"Brand notes: {style}\n"
        f"Allowed proofs:\n{json.dumps(brief.get('proofs', []), ensure_ascii=False)}\n\n"
        "Rewrite or remove any claim not supported by an allowed proof. "
        "Reading level grade 7 to 9. Keep one CTA per email. "
        "Do not add keys. Return JSON only."
    )
    return coerce_json(chat_json(client, system, user))

def validate_output(data: dict) -> List[str]:
    errors: List[str] = []
    emails = data.get("emails", [])
    if not isinstance(emails, list) or not emails:
        errors.append("emails missing")
    for i, e in enumerate(emails, 1):
        subj = (e.get("subject") or "").strip()
        body = (e.get("body") or "").strip()
        cta = (e.get("cta") or "").strip()
        if len(subj) > 50:
            errors.append(f"email {i} subject too long")
        if not cta:
            errors.append(f"email {i} missing CTA")
        if len(body.split()) > 160:
            errors.append(f"email {i} body too long")
    if not data.get("linkedin"):
        errors.append("linkedin posts missing")
    if not (isinstance(data.get("script"), str) and data.get("script").strip()):
        errors.append("script missing")
    return errors

def format_output_txt(data: dict) -> str:
    lines: List[str] = []
    lines.append("=== EMAIL SEQUENCE ===")
    for i, e in enumerate(data.get("emails", []), 1):
        subject = (e.get("subject") or "").strip()
        body = (e.get("body") or "").strip()
        cta = (e.get("cta") or "").strip()
        lines.append(f"\nStep {i}: {subject}")
        if body:
            lines.append(body)
        if cta:
            lines.append(f"CTA: {cta}")
        lines.append("-" * 40)

    lines.append("\n=== LINKEDIN POSTS ===")
    for i, p in enumerate(data.get("linkedin", []), 1):
        lines.append(f"\nPost {i}")
        lines.append((p or "").strip())
        lines.append("-" * 40)

    lines.append("\n=== TALK TRACK ===")
    lines.append((data.get("script") or "").strip())
    return "\n".join(lines)


# UI: header / hero / footer
def header_nav() -> None:
    st.markdown(
        """
<div class="bf-header">
  <div class="bf-left">
    <span class="bf-logo-badge">Deckflow</span>
  </div>
  <div class="bf-links">
    <a href="?page=presentation">Why Deckflow</a>
    <a href="?page=main">Product</a>
  </div>
  <div class="bf-right">
    <button class="bf-btn cta" onclick="window.parent.location.search='?page=main&scroll=agent'">Sign in</button>
  </div>
</div>
""",
        unsafe_allow_html=True,
    )

def hero_main() -> None:
    st.markdown(
        """
<div class="bf-hero">
  <div class="bf-chip">Built for GTM teams • Deck → Outbound in minutes</div>
  <h1 class="bf-h1">Turn Sales Presentations into Outbound Sequences</h1>
  <div class="bf-stats">
    <div class="bf-stat"><strong>6h → 1m</strong><span>Time to first draft</span></div>
    <div class="bf-stat"><strong>1 CTA</strong><span>Per email, enforced</span></div>
    <div class="bf-stat"><strong>Proof-tied</strong><span>Claims from your deck</span></div>
  </div>
  <ul class="bf-list">
    <li>⬇️ Upload your sales deck PDF or PPTX</li>
    <li>✉️ Get a 5 step email sequence with one clear CTA</li>
    <li>💬 Receive 3 LinkedIn posts ready to publish</li>
    <li>🎙️ Generate a 90 second talk track tied to your proofs</li>
  </ul>
</div>
""",
        unsafe_allow_html=True,
    )

    left, b1col, b2col, right = st.columns([1, 0.32, 0.32, 1])
    with b1col:
        if st.button("Try Now", type="primary", key="btn_try_now"):
            set_qs(page="main")
            st.session_state.scroll_to_agent = True
            st.rerun()
    with b2col:
        if st.button("See Why", key="btn_see_why"):
            set_qs(page="presentation")
            st.rerun()


def hero_presentation() -> None:
    st.markdown(
        """
<div class="bf-hero">
  <div class="bf-chip">Presentation</div>
  <h1 class="bf-h1">Why Deckflow, and how it helps your team</h1>
  <p class="bf-sub">Decision, value, and how the agent works end to end</p>
</div>
""",
        unsafe_allow_html=True,
    )
    if st.button("Back to Product", use_container_width=True, key="btn_back_to_product"):
        set_qs(page="main", scroll="agent")
        st.rerun()

def footer() -> None:
    st.markdown(
        """
<div class="bf-footer">
  <p>© 2025 Deckflow Mock | Demo only</p>
  <p><a href="?page=presentation">Why Deckflow</a> | <a href="?page=main">Product</a></p>
</div>
""",
        unsafe_allow_html=True,
    )


# Pages

def show_main_page() -> None:
    hero_main()

    st.markdown('<div id="agent"></div>', unsafe_allow_html=True)

    # smooth scroll triggers
    if st.session_state.scroll_to_agent:
        inject_scroll_to("agent")
        st.session_state.scroll_to_agent = False
    if get_qs("scroll", "") == "agent":
        inject_scroll_to("agent")
        # clear the param after use
        set_qs(page="main")

    # Step 1
    l1, c1, r1 = st.columns([1, 2, 1], gap="large")
    with c1:
        st.markdown('<h2 style="text-align:center;">Step 1. Upload your sales deck</h2>', unsafe_allow_html=True)
        uploaded = st.file_uploader("PDF or PPTX", type=["pdf", "pptx"])

    # Step 2
    l2, c2, r2 = st.columns([1, 2, 1], gap="large")
    with c2:
        st.markdown('<h2 style="text-align:center;">Step 2. Set persona and voice</h2>', unsafe_allow_html=True)
        persona = st.text_input("Target persona", "Head of Operations")
        tone = st.selectbox("Tone", ["Professional", "Friendly", "Direct", "Technical"])
        style = st.text_area("Brand voice notes", "Clear, concise, value focused.", height=80)
        run = st.button("Generate", type="primary", use_container_width=True, key="generate_btn")

    # Gate
    if st.session_state.final is None and not run:
        return
    if st.session_state.final is None and not uploaded:
        st.warning("Upload a deck first.")
        return

    # Heavy work (once)
    if run and st.session_state.final is None:
        with st.container():
            st.header("Step 3. We read your deck")
            tmp_path = Path(f"./_uploads/{uploaded.name}")
            tmp_path.parent.mkdir(parents=True, exist_ok=True)
            tmp_path.write_bytes(uploaded.read())

            with st.spinner("Extracting slides"):
                raw = extract_pdf_text(str(tmp_path)) if uploaded.name.lower().endswith(".pdf") else extract_pptx_text(str(tmp_path))
                if not raw.strip():
                    st.error("No text found in the deck. Add speaker notes or try another file.")
                    st.stop()
                outline = to_outline(raw)
                st.session_state.last_outline = outline

            with st.expander("Preview extracted outline", expanded=False):
                st.code(json.dumps(st.session_state.last_outline, ensure_ascii=False, indent=2), language="json")

        with st.container():
            st.header("Step 4. Draft your outbound pack")
            client = openai_client()

            with st.spinner("Planning brief"):
                brief = plan_brief(client, st.session_state.last_outline, persona, tone, style)

            with st.expander("Brief summary", expanded=True):
                st.write(f"Audience: {brief.get('audience','')}")
                pains = brief.get("pains", [])
                if pains:
                    st.write("Pains:")
                    for p in pains:
                        st.write(f"- {p}")
                proofs = brief.get("proofs", [])
                if proofs:
                    st.write("Top proofs:")
                    for p in proofs[:3]:
                        st.write(f"- {p}")

            with st.spinner("Generating drafts"):
                draft = generate_assets(client, brief)

            with st.spinner("Applying guardrails"):
                final = guardrail_review(client, draft, style, brief)

            issues = validate_output(final)
            if issues:
                st.info("Quality checks")
                for e in issues:
                    st.write(f"- {e}")

            st.session_state.final = final
            st.session_state.txt_content = format_output_txt(final)

    # Results (from cache)
    if st.session_state.final is not None:
        final = st.session_state.final

        with st.container():
            st.header("Results")
            tabs = st.tabs(["Emails", "LinkedIn", "Talk track", "JSON"])
            emails = final.get("emails", [])
            linkedin = final.get("linkedin", [])
            script = final.get("script", "")

            with tabs[0]:
                for i, e in enumerate(emails, 1):
                    st.markdown(f"**Step {i}: {e.get('subject','')}**")
                    st.write(e.get("body", ""))
                    st.write(f"CTA: {e.get('cta','')}")
                    st.divider()

            with tabs[1]:
                for i, p in enumerate(linkedin, 1):
                    st.markdown(f"**Post {i}**")
                    st.write(p)
                    st.divider()

            with tabs[2]:
                st.markdown("**Talk track**")
                st.write(script)

            with tabs[3]:
                st.code(json.dumps(final, ensure_ascii=False, indent=2), language="json")

            c1, c2 = st.columns(2)
            with c1:
                st.download_button(
                    "Download TXT",
                    data=st.session_state.txt_content,
                    file_name="outbound.txt",
                    mime="text/plain",
                    use_container_width=True,
                    key="dl_txt",
                )
            with c2:
                st.download_button(
                    "Download JSON",
                    data=json.dumps(final, ensure_ascii=False, indent=2),
                    file_name="outbound.json",
                    mime="application/json",
                    use_container_width=True,
                    key="dl_json",
                )

            st.button(
                "Reset Result",
                on_click=lambda: st.session_state.update({"final": None, "txt_content": "", "last_outline": None})
            )

def show_presentation_page() -> None:
    hero_presentation()

    with st.container():
        st.header("Problem")
        st.write(
            "- Sales decks hold the best proof and messaging, but marketing rewrites them by hand.\n"
            "- Outbound sequences take hours each week.\n"
            "- Messages drift from the source deck and lose accuracy."
        )

    with st.container():
        st.header("Solution")
        st.write(
            "Deckflow turns a sales presentation into:\n"
            "- A 5 step email sequence with one clear CTA per email\n"
            "- Three LinkedIn posts\n"
            "- A 90 second talk track\n"
            "All claims map to proofs extracted from the deck."
        )

    with st.container():
        st.header("How this accelerates the team")
        st.write(
            "- First draft in minutes, not hours\n"
            "- Consistent voice and claims from the same source\n"
            "- Guardrails remove unsupported claims\n"
            "- Reusable TXT and JSON exports"
        )

    with st.container():
        st.header("Return on investment")
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Time saved", "4–6 hrs/wk", "per marketer")
        with col2:
            st.metric("First draft", "≈ 1 min", "from upload")
        with col3:
            st.metric("Edit cycles", "-50%", "fewer rounds")
        st.write(
        "Example: a team of 3 marketers saves 12–18 hours a week. "
        "At \$80/hour, that’s \$960–\$1,440 in weekly value, plus faster launches."
        )
        st.markdown("</div>", unsafe_allow_html=True)


    with st.container():
        st.header("Clean and simple UI/UX")
        st.write(
            "- One page to upload and set voice\n"
            "- Tabs for emails, LinkedIn, and talk track\n"
            "- Download TXT for sharing, JSON for systems"
        )

    with st.container():
        st.header("Flow")
        st.write(
            "- Upload a PDF or PPTX\n"
            "- Pick persona and tone, add brand notes\n"
            "- The agent extracts proofs and builds a brief\n"
            "- The agent generates drafts and applies guardrails\n"
            "- You review, export TXT/JSON, and ship"
        )

    with st.container():
        st.header("Architecture")
        st.write(
            "- Extract with pdfplumber and python-pptx\n"
            "- Plan brief from deck outline\n"
            "- Generate emails, posts, and script\n"
            "- Guardrails enforce proofs and tone\n"
            "- Validate length and structure"
        )

    with st.container():
        st.header("Decision making")
        st.write(
            "- Chose deck → outbound for clear value and speed to impact\n"
            "- JSON + TXT outputs serve both humans and systems\n"
            "- Kept UI focused to reduce friction\n"
            "- Guardrails reduce legal and compliance risk"
        )

    with st.container():
        st.header("Roadmap")
        st.write(
            "- Subject line A/B suggestions\n"
            "- CRM export helpers (HubSpot, Salesforce blocks)\n"
            "- Team presets for voice and CTAs\n"
            "- PDF export with brand header and footer"
        )


# App entry

def main() -> None:
    st.set_page_config(page_title="Deckflow", page_icon="📨", layout="centered")
    load_css("styles.css")
    header_nav()

    page = get_qs("page", "main")
    if page == "presentation":
        show_presentation_page()
    else:
        show_main_page()

    footer()

if __name__ == "__main__":
    main()
